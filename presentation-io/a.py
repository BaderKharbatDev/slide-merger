import zipfile
from pptx import Presentation     
import os, shutil

def openPresentation(filepath, is_main_file):    
    with zipfile.ZipFile(filepath, 'r') as zip_ref:
        if is_main_file:
            zip_ref.extractall("./files/main")
        else:
            zip_ref.extractall("./files/sub-presentations")


def createBlankMainPresentation():
    #deletes everything in the folder from prior testing
    deleteMainProgramFile()

    # creates a blank pres file
    root = Presentation()
    root.save("./program-files/current.pptx")

    #unzips it
    with zipfile.ZipFile("./program-files/current.pptx", 'r') as zip_ref:
        zip_ref.extractall("./program-files/main")

    #removes the original pptx file
    os.remove("program-files/current.pptx")

def deleteMainProgramFile():
    dir = 'program-files/main/'
    for files in os.listdir(dir):
        path = os.path.join(dir, files)
        try:
            shutil.rmtree(path)
        except OSError:
            os.remove(path)

createBlankMainPresentation()


