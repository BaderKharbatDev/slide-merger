from pptx import Presentation

prs = Presentation('../Test-Files/original.pptx')
print(len(prs.slides))
i = 0
for slide in prs.slides: 
    if i % 2 == 0:
        prs.slides.pop(slide)
    print(slide)
    i+=1
prs.save('../Test-Files/new.pptx')